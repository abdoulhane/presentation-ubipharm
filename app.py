
import asyncio, hashlib, io, re, shutil, subprocess, tempfile, zipfile
from collections import Counter
from pathlib import Path

import edge_tts
import fitz
import streamlit as st
from groq import Groq
from pptx import Presentation

IGNORED = [
    r"^\s*Département\s+douane\s*$",
    r"^\s*Date\s*$",
    r"^\s*Titre de la présentation.*Émetteur\s*$",
    r"^\s*\d+\s*$",
]

EDGE_VOICES = {
    "Denise — femme, France": "fr-FR-DeniseNeural",
    "Henri — homme, France": "fr-FR-HenriNeural",
    "Eloise — femme, France": "fr-FR-EloiseNeural",
    "Remy — homme, France": "fr-FR-RemyMultilingualNeural",
    "Vivienne — femme, France": "fr-FR-VivienneMultilingualNeural",
}

def norm(t):
    t = t.replace("\u00a0", " ")
    t = re.sub(r"[ \t]+", " ", t)
    return re.sub(r"\n{3,}", "\n\n", t).strip()

def cmp_norm(t):
    return re.sub(r"\s+", " ", norm(t).lower()).strip(" .,:;-")

def shape_text(shape):
    out = []
    if hasattr(shape, "shapes"):
        for s in shape.shapes:
            out += shape_text(s)
    if getattr(shape, "has_text_frame", False):
        txt = "\n".join(norm(p.text) for p in shape.text_frame.paragraphs if norm(p.text))
        if txt:
            out.append(txt)
    if getattr(shape, "has_table", False):
        rows = []
        for row in shape.table.rows:
            vals = [norm(c.text) for c in row.cells if norm(c.text)]
            if vals:
                rows.append(" | ".join(vals))
        if rows:
            out.append("\n".join(rows))
    return out

def extract_slides(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        blocks, seen = [], set()
        for sh in slide.shapes:
            for b in shape_text(sh):
                k = cmp_norm(b)
                if k and k not in seen:
                    seen.add(k)
                    blocks.append(norm(b))
        slides.append({"number": i, "raw_blocks": blocks})
    return slides

def clean_slides(slides, custom_ignored):
    counts = Counter()
    for s in slides:
        for b in set(cmp_norm(x) for x in s["raw_blocks"] if len(cmp_norm(x)) >= 70):
            counts[b] += 1
    repeated = {k for k, v in counts.items() if v >= 3}
    cleaned = []
    for s in slides:
        kept = []
        for b in s["raw_blocks"]:
            if any(re.match(p, b, re.I) for p in IGNORED):
                continue
            if any(x.lower() in b.lower() for x in custom_ignored if x):
                continue
            if cmp_norm(b) in repeated:
                continue
            kept.append(b)
        cleaned.append({"number": s["number"], "clean_text": "\n\n".join(kept).strip()})
    return cleaned

def get_groq_key():
    try:
        return st.secrets["GROQ_API_KEY"]
    except Exception:
        return ""

def narration_prompt(slide, previous, target_seconds):
    return f"""
    Rédige uniquement le texte oral final en français.
    Fais une phrase d'introduction global pour la présentation
    RÈGLES :
    - Utilise exclusivement les informations présentes dans CONTENU.
    - N'invente aucune information pour atteindre une longueur donnée.
    - Si le contenu est très court ou correspond à une page de titre,
    une narration de 1 ou 2 phrases suffit.
    - Ne développe pas la signification d'un terme si elle n'est pas
    explicitement donnée dans le contenu.
    - Ne dis jamais "slide" ou "diapositive".
    - Ton naturel, professionnel et simple.
    - Retourne uniquement ce qui doit être prononcé.

CONTENU :
{slide["clean_text"]}

NARRATION PRÉCÉDENTE :
{previous[-1000:] if previous else "[Aucune]"}
""".strip()

def generate_narration(key, model, prompt):
    client = Groq(api_key=key)
    r = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": "Tu rédiges des narrations professionnelles fidèles aux sources."},
            {"role": "user", "content": prompt},
        ],
        temperature=0.3,
        max_tokens=1000,
    )
    return norm(r.choices[0].message.content or "")

async def edge_audio_async(text, voice, rate, pitch):
    c = edge_tts.Communicate(text=text, voice=voice, rate=rate, pitch=pitch)
    data = bytearray()
    async for chunk in c.stream():
        if chunk["type"] == "audio":
            data.extend(chunk["data"])
    if not data:
        raise RuntimeError("Edge TTS n'a produit aucun audio.")
    return bytes(data)

def edge_audio(text, voice, rate, pitch):
    return asyncio.run(edge_audio_async(text, voice, rate, pitch))

def find_cmd(names):
    for n in names:
        p = shutil.which(n)
        if p:
            return p
    return None

def mp3_to_wav(mp3):
    ffmpeg = find_cmd(["ffmpeg"])
    if not ffmpeg:
        raise RuntimeError("FFmpeg introuvable.")
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        a, w = td/"a.mp3", td/"a.wav"
        a.write_bytes(mp3)
        r = subprocess.run([ffmpeg, "-y", "-i", str(a), "-ar", "24000", "-ac", "1", str(w)],
                           capture_output=True, text=True, timeout=120)
        if r.returncode != 0 or not w.exists():
            raise RuntimeError(r.stderr[-1500:])
        return w.read_bytes()

def silent_wav(seconds=2):
    import wave
    buf = io.BytesIO()
    rate = 24000
    with wave.open(buf, "wb") as wf:
        wf.setnchannels(1); wf.setsampwidth(2); wf.setframerate(rate)
        wf.writeframes(b"\x00\x00" * int(rate * seconds))
    return buf.getvalue()

def render_slides(powerpoint_bytes, workdir, extension=".pptx"):
    soffice = find_cmd(["libreoffice", "soffice"])
    if not soffice:
        raise RuntimeError("LibreOffice introuvable.")

    extension = extension.lower()
    if extension not in [".pptx", ".pptm"]:
        raise RuntimeError(f"Format PowerPoint non pris en charge : {extension}")

    powerpoint_file = workdir / f"presentation{extension}"
    powerpoint_file.write_bytes(powerpoint_bytes)

    r = subprocess.run(
        [
            soffice,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(workdir),
            str(powerpoint_file),
        ],
        capture_output=True,
        text=True,
        timeout=180,
    )
    pdf = workdir/"presentation.pdf"
    if r.returncode != 0 or not pdf.exists():
        raise RuntimeError("Conversion PowerPoint → PDF impossible.")
    doc = fitz.open(pdf)
    out = []
    for i, page in enumerate(doc, 1):
        pix = page.get_pixmap(matrix=fitz.Matrix(2,2), alpha=False)
        p = workdir/f"slide_{i:03d}.png"
        pix.save(p)
        out.append(p)
    doc.close()
    return out

def build_video(powerpoint_bytes, audios, slide_count, resolution="1280x720", extension=".pptx"):
    ffmpeg = find_cmd(["ffmpeg"])
    if not ffmpeg:
        raise RuntimeError("FFmpeg introuvable.")
    w, h = map(int, resolution.split("x"))
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        pngs = render_slides(powerpoint_bytes, td, extension)
        segs = []
        for i, png in enumerate(pngs, 1):
            wav = td/f"a_{i}.wav"
            wav.write_bytes(audios.get(i, silent_wav()))
            seg = td/f"s_{i}.mp4"
            vf = f"scale={w}:{h}:force_original_aspect_ratio=decrease,pad={w}:{h}:(ow-iw)/2:(oh-ih)/2,setsar=1"
            r = subprocess.run([ffmpeg,"-y","-loop","1","-framerate","25","-i",str(png),"-i",str(wav),
                                "-vf",vf,"-c:v","libx264","-preset","veryfast","-pix_fmt","yuv420p",
                                "-c:a","aac","-shortest",str(seg)],
                               capture_output=True,text=True,timeout=240)
            if r.returncode != 0:
                raise RuntimeError(r.stderr[-1500:])
            segs.append(seg)
        concat = td/"concat.txt"
        concat.write_text("\n".join(f"file '{p.as_posix()}'" for p in segs), encoding="utf-8")
        out = td/"presentation_narree.mp4"
        r = subprocess.run([ffmpeg,"-y","-f","concat","-safe","0","-i",str(concat),"-c","copy",str(out)],
                           capture_output=True,text=True,timeout=240)
        if r.returncode != 0 or not out.exists():
            raise RuntimeError(r.stderr[-1500:])
        return out.read_bytes()

st.set_page_config(page_title="Présentation IA V5", page_icon="🎬", layout="wide")
st.title("🎬 Présentation IA — V5")
st.caption("Groq + Edge TTS + génération MP4 • formats .pptx et .pptm")

with st.sidebar:
    key = get_groq_key() or st.text_input("Clé API Groq", type="password")
    model = "openai/gpt-oss-safeguard-20b"

    st.sidebar.write(
        "Modèle de narration : openai/gpt-oss-safeguard-20b"
    )
    target_seconds = st.slider("Durée cible par slide", 15, 90, 40, 5)
    ignored_text = st.text_area("Expressions à ignorer", "Département douane\nTitre de la présentation\nÉmetteur")
    ignored = [x.strip() for x in ignored_text.splitlines() if x.strip()]
    voice_label = st.selectbox("Voix française", list(EDGE_VOICES.keys()))
    voice = EDGE_VOICES[voice_label]
    rate_i = st.slider("Vitesse (%)", -30, 30, -5, 5)
    pitch_i = st.slider("Hauteur (Hz)", -20, 20, 0, 5)
    rate, pitch = f"{rate_i:+d}%", f"{pitch_i:+d}Hz"
    resolution = st.selectbox("Résolution", ["1280x720", "1920x1080"])

uploaded = st.file_uploader("Dépose ton PowerPoint", type=["pptx", "pptm"])
if not uploaded:
    st.stop()

powerpoint_bytes = uploaded.getvalue()
extension = Path(uploaded.name).suffix.lower()
file_hash = hashlib.md5(powerpoint_bytes).hexdigest()

if extension not in [".pptx", ".pptm"]:
    st.error("Format non pris en charge. Utilise un fichier .pptx ou .pptm.")
    st.stop()

try:
    slides = clean_slides(extract_slides(powerpoint_bytes), ignored)
except Exception as exc:
    st.error(
        "Impossible de lire cette présentation. "
        "Certains fichiers .pptm très particuliers peuvent poser problème.\n\n"
        f"Détail : {exc}"
    )
    st.stop()

if st.session_state.get("file_hash") != file_hash:
    st.session_state["file_hash"] = file_hash
    st.session_state["narrations"] = {}
    st.session_state["audios"] = {}
    st.session_state["video"] = None

narr = st.session_state["narrations"]
audios = st.session_state["audios"]
tabs = st.tabs(["1. Contenu", "2. Narrations", "3. Voix", "4. Vidéo"])

with tabs[0]:
    for s in slides:
        with st.expander(f"Slide {s['number']}"):
            st.text_area("Texte", s["clean_text"], height=140, disabled=True, key=f"src{s['number']}")

with tabs[1]:
    if st.button("✨ Générer toutes les narrations", type="primary", disabled=not key):
        prev = ""
        bar = st.progress(0)
        for i, s in enumerate(slides):
            if s["clean_text"]:
                narr[s["number"]] = generate_narration(key, model, narration_prompt(s, prev, target_seconds))
                prev = narr[s["number"]]
            bar.progress((i+1)/len(slides))
        st.session_state["narrations"] = narr
        st.success("Narrations générées.")
    for s in slides:
        n = s["number"]
        cur = narr.get(n, "")
        edit = st.text_area(f"Slide {n}", cur, height=130, key=f"narr{n}_{hash(cur)}")
        if edit != cur:
            narr[n] = edit
            audios.pop(n, None)
            st.session_state["video"] = None

with tabs[2]:
    st.info("Edge TTS ne charge aucun modèle lourd en mémoire.")
    if st.button("🔊 Générer tous les audios", type="primary"):
        bar = st.progress(0)
        for i, s in enumerate(slides):
            n = s["number"]
            text = narr.get(n, "").strip()
            if text:
                audios[n] = mp3_to_wav(edge_audio(text, voice, rate, pitch))
            bar.progress((i+1)/len(slides))
        st.session_state["audios"] = audios
        st.success("Audios générés.")
    for s in slides:
        if s["number"] in audios:
            st.audio(audios[s["number"]], format="audio/wav")

with tabs[3]:
    libreoffice_ok = bool(find_cmd(["libreoffice", "soffice"]))
    ffmpeg_ok = bool(find_cmd(["ffmpeg"]))
    st.write("LibreOffice :", "✅" if libreoffice_ok else "❌")
    st.write("FFmpeg :", "✅" if ffmpeg_ok else "❌")
    if st.button("🎬 Générer le MP4", type="primary", disabled=not (libreoffice_ok and ffmpeg_ok)):
        with st.spinner("Création de la vidéo…"):
            st.session_state["video"] = build_video(powerpoint_bytes, audios, len(slides), resolution, extension)
    if st.session_state.get("video"):
        st.video(st.session_state["video"])
        st.download_button("⬇️ Télécharger la vidéo", st.session_state["video"], "presentation_narree.mp4", "video/mp4")
